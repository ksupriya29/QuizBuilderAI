from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
from docx import Document
from pypdf import PdfReader
from openpyxl import load_workbook
import requests, os, io, json, csv
from dotenv import load_dotenv

load_dotenv()
app = Flask(__name__)
CORS(app)

DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")

SUPPORTED_EXTENSIONS = {
    ".pptx": "PowerPoint",
    ".ppt":  "PowerPoint (old format)",
    ".docx": "Word Document",
    ".doc":  "Word Document (old format)",
    ".pdf":  "PDF Document",
    ".txt":  "Text File",
    ".csv":  "CSV File",
    ".xlsx": "Excel File",
}

def extract_text_from_pptx(file_bytes):
    """Extract text from .pptx files"""
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        text = " ".join(
            shape.text for shape in slide.shapes if hasattr(shape, "text")
        )
        if text.strip():
            slides.append({"slide": i + 1, "text": text.strip()})
    return slides

def extract_text_from_docx(file_bytes):
    """Extract text from .docx files"""
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Group into chunks (simulate slides/pages)
    chunks = []
    chunk_size = max(1, len(paragraphs) // 10) if len(paragraphs) > 10 else len(paragraphs)
    for i in range(0, len(paragraphs), chunk_size):
        chunk_text = "\n".join(paragraphs[i:i+chunk_size])
        if chunk_text.strip():
            chunks.append({"slide": len(chunks) + 1, "text": chunk_text.strip()})
    if not chunks:
        chunks.append({"slide": 1, "text": "\n".join(paragraphs)})
    return chunks

def extract_text_from_pdf(file_bytes):
    """Extract text from PDF files"""
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text.strip():
            pages.append({"slide": i + 1, "text": text.strip()})
    return pages

def extract_text_from_txt(file_bytes):
    """Extract text from plain text files"""
    text = file_bytes.decode("utf-8", errors="replace")
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    # Group into chunks
    chunks = []
    chunk_size = max(1, len(lines) // 10) if len(lines) > 10 else len(lines)
    for i in range(0, len(lines), chunk_size):
        chunk_text = "\n".join(lines[i:i+chunk_size])
        if chunk_text.strip():
            chunks.append({"slide": len(chunks) + 1, "text": chunk_text.strip()})
    if not chunks:
        chunks.append({"slide": 1, "text": text.strip()})
    return chunks

def extract_text_from_csv(file_bytes):
    """Extract text from CSV files"""
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        rows.append(" | ".join(row))
    chunks = []
    chunk_size = max(1, len(rows) // 10) if len(rows) > 10 else len(rows)
    for i in range(0, len(rows), chunk_size):
        chunk_text = "\n".join(rows[i:i+chunk_size])
        if chunk_text.strip():
            chunks.append({"slide": len(chunks) + 1, "text": chunk_text.strip()})
    if not chunks:
        chunks.append({"slide": 1, "text": "\n".join(rows)})
    return chunks

def extract_text_from_xlsx(file_bytes):
    """Extract text from Excel files"""
    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    rows = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                rows.append(row_text)
    wb.close()
    chunks = []
    chunk_size = max(1, len(rows) // 10) if len(rows) > 10 else len(rows)
    for i in range(0, len(rows), chunk_size):
        chunk_text = "\n".join(rows[i:i+chunk_size])
        if chunk_text.strip():
            chunks.append({"slide": len(chunks) + 1, "text": chunk_text.strip()})
    if not chunks:
        chunks.append({"slide": 1, "text": "\n".join(rows)})
    return chunks

def extract_text_from_ppt_old(file_bytes):
    """
    Try to extract text from old .ppt format using pywin32 (Windows only, requires PowerPoint installed).
    Falls back gracefully if not available.
    """
    try:
        import win32com.client
        import tempfile
        import pythoncom

        pythoncom.CoInitialize()
        temp_path = None
        try:
            # Write to temp file
            with tempfile.NamedTemporaryFile(suffix=".ppt", delete=False) as tmp:
                tmp.write(file_bytes)
                temp_path = tmp.name

            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = False
            presentation = powerpoint.Presentations.Open(temp_path, WithWindow=False)
            
            slides = []
            for i, slide in enumerate(presentation.Slides):
                text_parts = []
                for shape in slide.Shapes:
                    if shape.HasTextFrame:
                        text_parts.append(shape.TextFrame.TextRange.Text)
                text = " ".join(text_parts)
                if text.strip():
                    slides.append({"slide": i + 1, "text": text.strip()})
            
            presentation.Close()
            return slides
        except Exception as e:
            print(f"[WARN] pywin32 extraction failed: {e}")
            return None
        finally:
            if temp_path and os.path.exists(temp_path):
                os.unlink(temp_path)
            try:
                pythoncom.CoUninitialize()
            except:
                pass
    except ImportError:
        print("[WARN] pywin32 not installed, cannot parse .ppt files")
        return None

def extract_text_from_doc_old(file_bytes):
    """
    Try to extract text from old .doc format using pywin32 (Windows only, requires Word installed).
    Falls back gracefully if not available.
    """
    try:
        import win32com.client
        import tempfile
        import pythoncom

        pythoncom.CoInitialize()
        temp_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
                tmp.write(file_bytes)
                temp_path = tmp.name

            word = win32com.client.Dispatch("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(temp_path)
            
            paragraphs = [p.text for p in doc.Paragraphs if p.text.strip()]
            
            doc.Close()
            
            chunks = []
            chunk_size = max(1, len(paragraphs) // 10) if len(paragraphs) > 10 else len(paragraphs)
            for i in range(0, len(paragraphs), chunk_size):
                chunk_text = "\n".join(paragraphs[i:i+chunk_size])
                if chunk_text.strip():
                    chunks.append({"slide": len(chunks) + 1, "text": chunk_text.strip()})
            if not chunks:
                chunks.append({"slide": 1, "text": "\n".join(paragraphs)})
            return chunks
        except Exception as e:
            print(f"[WARN] pywin32 .doc extraction failed: {e}")
            return None
        finally:
            if temp_path and os.path.exists(temp_path):
                os.unlink(temp_path)
            try:
                pythoncom.CoUninitialize()
            except:
                pass
    except ImportError:
        print("[WARN] pywin32 not installed, cannot parse .doc files")
        return None

@app.route("/parse", methods=["POST"])
def parse_file():
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    
    file = request.files["file"]
    filename = file.filename.lower()
    file_bytes = file.read()
    
    # Determine file extension
    ext = os.path.splitext(filename)[1].lower()
    
    if ext not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(SUPPORTED_EXTENSIONS.keys())
        return jsonify({
            "error": f"Unsupported file type '{ext}'. Supported types: {supported}"
        }), 400
    
    # Extract text based on file type
    slides = None
    
    if ext == ".pptx":
        slides = extract_text_from_pptx(file_bytes)
    elif ext == ".ppt":
        slides = extract_text_from_ppt_old(file_bytes)
        if slides is None:
            return jsonify({
                "error": "Could not parse .ppt file. Please convert to .pptx format and try again."
            }), 400
    elif ext == ".docx":
        slides = extract_text_from_docx(file_bytes)
    elif ext == ".doc":
        slides = extract_text_from_doc_old(file_bytes)
        if slides is None:
            return jsonify({
                "error": "Could not parse .doc file. Please convert to .docx format and try again."
            }), 400
    elif ext == ".pdf":
        slides = extract_text_from_pdf(file_bytes)
    elif ext == ".txt":
        slides = extract_text_from_txt(file_bytes)
    elif ext == ".csv":
        slides = extract_text_from_csv(file_bytes)
    elif ext == ".xlsx":
        slides = extract_text_from_xlsx(file_bytes)
    
    if not slides:
        return jsonify({
            "error": "No text content could be extracted from the file."
        }), 400
    
    total_words = sum(len(s["text"].split()) for s in slides)
    
    return jsonify({
        "slide_count": len(slides),
        "word_count": total_words,
        "slides": slides,
        "file_type": SUPPORTED_EXTENSIONS.get(ext, ext),
        "filename": file.filename
    })

@app.route("/generate", methods=["POST"])
def generate_quiz():
    try:
        data = request.json
        slide_text = data["slide_text"]
        num_questions = data["num_questions"]
        difficulty = data["difficulty"]

        prompt = f"""You are a quiz generator. Given the slide content below, generate exactly {num_questions} MCQ questions at {difficulty} difficulty.

Rules:
- Each question has exactly 4 options (A, B, C, D)
- Only one correct answer per question
- Distractors must be plausible
- Include explanation for each wrong option

Return ONLY valid JSON array (no markdown, no extra text):
[
  {{
    "question": "...",
    "options": {{"A": "...", "B": "...", "C": "...", "D": "..."}},
    "correct": "B",
    "explanations": {{
      "A": "Why A is wrong",
      "C": "Why C is wrong",
      "D": "Why D is wrong"
    }}
  }}
]

Slide Content:
{slide_text}"""

        if DEEPSEEK_API_KEY and DEEPSEEK_API_KEY.startswith("sk-or-"):
            api_url = "https://openrouter.ai/api/v1/chat/completions"
            model = "deepseek/deepseek-chat"
            headers = {
                "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
                "Content-Type": "application/json",
            }
        else:
            api_url = "https://api.deepseek.com/chat/completions"
            model = "deepseek-chat"
            headers = {
                "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
                "Content-Type": "application/json"
            }

        # Retry up to 2 times on empty/transient errors
        retries = 0
        max_retries = 2
        raw = None
        while retries <= max_retries:
            response = requests.post(api_url, headers=headers, json={
                "model": model,
                "messages": [{"role": "user", "content": prompt}],
                "temperature": 0.7
            })

            if not response.ok:
                error_detail = response.text
                print(f"[ERROR] API returned {response.status_code}: {error_detail}")
                return jsonify({"error": f"API error {response.status_code}: {error_detail}"}), response.status_code

            raw = response.text.strip()
            if raw:
                break
            retries += 1
            print(f"[WARN] Empty API response, retry {retries}/{max_retries}")
        
        if not raw:
            return jsonify({"error": "AI API returned empty response after retries"}), 502

        # Try to parse the response - sometimes the API wrapper sends whitespace before JSON
        # or the content itself has escaped JSON that needs double-parsing
        print(f"[DEBUG] Raw response length: {len(raw)}, first 200: {raw[:200]}")
        
        result = response.json()
        content = result["choices"][0]["message"]["content"]
        print(f"[DEBUG] Content from API (first 300): {content[:300]}")
        
        # Clean the content - remove markdown code fences if present
        content = content.replace("```json", "").replace("```", "").strip()
        
        # If content is empty, try direct content from raw response
        if not content:
            print(f"[WARN] Empty content field, trying raw response body")
            return jsonify({"error": "AI returned empty content"}), 502
        
        # Try to parse as JSON - if content is double-escaped (string inside string), handle it
        try:
            questions = json.loads(content)
        except json.JSONDecodeError:
            # Content might be double-escaped JSON string - try to extract
            print(f"[WARN] Direct parse failed, content starts with: {content[:100]}")
            # Try to find JSON array brackets in the content
            start = content.find('[')
            end = content.rfind(']')
            if start != -1 and end != -1:
                extracted = content[start:end+1]
                print(f"[DEBUG] Extracted JSON substring: {extracted[:200]}")
                questions = json.loads(extracted)
            else:
                raise
        
        return jsonify({"questions": questions})
    except Exception as e:
        print(f"[ERROR] Exception in generate_quiz: {e}")
        return jsonify({"error": f"Server error: {str(e)}"}), 500

if __name__ == "__main__":
    app.run(debug=True, port=5000)